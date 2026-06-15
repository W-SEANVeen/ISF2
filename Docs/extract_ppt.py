import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
import glob, os

files = [f for f in glob.glob('Docs/*.pptx') if not f.startswith('Docs/~')]
files.sort(key=os.path.getmtime, reverse=True)
print(f"Found files: {files}")
pptx_path = files[0] if files else None
print(f"Using: {pptx_path}")

prs = Presentation(pptx_path)

print(f'Total slides: {len(prs.slides)}')
for i, slide in enumerate(prs.slides, 1):
    layout = slide.slide_layout
    print(f'=== SLIDE {i} (layout: {layout.name}) ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(f'  [{shape.name}] {shape.text.strip()[:500]}')
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                print(f'  TABLE: {" | ".join(cells)}')
    print()
